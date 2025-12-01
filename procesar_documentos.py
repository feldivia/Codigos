"""
Módulo de Procesamiento de Documentos
=====================================
Funcionalidades para extraer información estructurada de documentos PDF, Word y PowerPoint.

Librerías utilizadas:
- PDF: pdfplumber (extracción de texto y tablas), pypdf (metadatos)
- Word (.docx): python-docx (contenido), pandoc (conversión markdown)
- PowerPoint (.pptx): python-pptx (slides y contenido)

Autor: Sistema de Propuestas Automatizadas
"""

import os
import re
import json
import tempfile
from typing import Dict, List, Optional, Tuple, Any
from dataclasses import dataclass, field, asdict
from enum import Enum
from pathlib import Path
import subprocess

# Importaciones para PDF
try:
    import pdfplumber
    PDF_AVAILABLE = True
except ImportError:
    PDF_AVAILABLE = False

try:
    from pypdf import PdfReader
    PYPDF_AVAILABLE = True
except ImportError:
    PYPDF_AVAILABLE = False

# Importaciones para Word
try:
    from docx import Document as DocxDocument
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False

# Importaciones para PowerPoint
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

# Para procesamiento de lenguaje natural (extracción de keywords)
try:
    from collections import Counter
    import string
except ImportError:
    pass


class DocumentType(Enum):
    """Tipos de documentos soportados"""
    PDF = "pdf"
    WORD = "docx"
    POWERPOINT = "pptx"
    UNKNOWN = "unknown"


class ProcessingStatus(Enum):
    """Estados de procesamiento"""
    SUCCESS = "success"
    ERROR = "error"
    PENDING = "pending"


@dataclass
class DocumentMetadata:
    """Metadatos extraídos de un documento"""
    filename: str = ""
    file_type: str = ""
    title: str = ""
    author: str = ""
    creation_date: str = ""
    page_count: int = 0
    word_count: int = 0


@dataclass
class ExtractedContent:
    """Contenido extraído de un documento"""
    # Información general
    titulo_propuesta: str = ""
    palabras_claves: List[str] = field(default_factory=list)
    cliente: str = ""
    industria: str = ""
    
    # Contexto
    problema: str = ""
    objetivo_general: str = ""
    objetivos_secundarios: List[str] = field(default_factory=list)
    
    # Alcance
    alcance_funcional: str = ""
    alcance_tecnico: str = ""
    alcance_geografico: str = ""
    limitaciones: str = ""
    
    # Metadatos
    metadata: DocumentMetadata = field(default_factory=DocumentMetadata)
    raw_text: str = ""


@dataclass
class ProcessingResult:
    """Resultado del procesamiento de un documento"""
    filename: str
    status: ProcessingStatus
    content: Optional[ExtractedContent] = None
    error_message: str = ""
    
    def to_dict(self) -> Dict:
        """Convierte el resultado a diccionario"""
        return {
            "filename": self.filename,
            "status": self.status.value,
            "content": asdict(self.content) if self.content else None,
            "error_message": self.error_message
        }


class DocumentProcessor:
    """
    Clase principal para procesar documentos y extraer información estructurada.
    
    Soporta:
    - PDF (.pdf)
    - Word (.docx)
    - PowerPoint (.pptx)
    """
    
    # Palabras vacías en español para filtrar keywords
    STOPWORDS_ES = {
        'el', 'la', 'los', 'las', 'un', 'una', 'unos', 'unas', 'de', 'del', 'al',
        'a', 'ante', 'bajo', 'cabe', 'con', 'contra', 'desde', 'en', 'entre',
        'hacia', 'hasta', 'para', 'por', 'según', 'sin', 'so', 'sobre', 'tras',
        'y', 'e', 'ni', 'o', 'u', 'pero', 'sino', 'que', 'si', 'como', 'cuando',
        'donde', 'mientras', 'aunque', 'porque', 'pues', 'ya', 'ser', 'estar',
        'haber', 'tener', 'hacer', 'poder', 'decir', 'ir', 'ver', 'dar', 'saber',
        'querer', 'llegar', 'pasar', 'deber', 'poner', 'parecer', 'quedar',
        'creer', 'hablar', 'llevar', 'dejar', 'seguir', 'encontrar', 'llamar',
        'venir', 'pensar', 'salir', 'volver', 'tomar', 'conocer', 'vivir',
        'sentir', 'tratar', 'mirar', 'contar', 'empezar', 'esperar', 'buscar',
        'existir', 'entrar', 'trabajar', 'escribir', 'perder', 'producir',
        'ocurrir', 'entender', 'pedir', 'recibir', 'recordar', 'terminar',
        'permitir', 'aparecer', 'conseguir', 'comenzar', 'servir', 'sacar',
        'necesitar', 'mantener', 'resultar', 'leer', 'caer', 'cambiar', 'presentar',
        'crear', 'abrir', 'considerar', 'oír', 'acabar', 'convertir', 'ganar',
        'formar', 'traer', 'partir', 'morir', 'aceptar', 'realizar', 'suponer',
        'comprender', 'lograr', 'explicar', 'preguntar', 'tocar', 'reconocer',
        'esto', 'esta', 'estos', 'estas', 'ese', 'esa', 'esos', 'esas', 'aquel',
        'aquella', 'aquellos', 'aquellas', 'yo', 'tú', 'él', 'ella', 'nosotros',
        'vosotros', 'ellos', 'ellas', 'me', 'te', 'se', 'nos', 'os', 'le', 'les',
        'lo', 'la', 'los', 'las', 'mí', 'ti', 'sí', 'conmigo', 'contigo', 'consigo',
        'mi', 'tu', 'su', 'nuestro', 'vuestro', 'mío', 'tuyo', 'suyo', 'más',
        'menos', 'muy', 'mucho', 'poco', 'todo', 'nada', 'algo', 'cada', 'otro',
        'mismo', 'tanto', 'tan', 'así', 'bien', 'mal', 'mejor', 'peor', 'sólo',
        'solo', 'siempre', 'nunca', 'también', 'además', 'ahora', 'después',
        'antes', 'entonces', 'luego', 'aquí', 'allí', 'así', 'etc', 'puede',
        'pueden', 'debe', 'deben', 'será', 'serán', 'sido', 'siendo', 'son',
        'fue', 'fueron', 'era', 'eran', 'sea', 'sean', 'tiene', 'tienen', 'tenía',
        'han', 'ha', 'había', 'hemos', 'hay', 'cual', 'cuales', 'quien', 'quienes',
        'cuyo', 'cuya', 'cuyos', 'cuyas', 'qué', 'cuál', 'cuáles', 'quién',
        'quiénes', 'dónde', 'cuándo', 'cuánto', 'cuánta', 'cuántos', 'cuántas',
        'cómo', 'www', 'http', 'https', 'com', 'org', 'net', 'page', 'slide'
    }
    
    # Patrones para identificar secciones
    SECTION_PATTERNS = {
        'titulo': [
            r'(?:título|titulo|nombre)[\s:]+(.+)',
            r'^#\s*(.+)$',
            r'propuesta[\s:]+(.+)',
        ],
        'cliente': [
            r'(?:cliente|customer|empresa)[\s:]+(.+)',
            r'(?:para|for)[\s:]+(.+?)(?:\n|$)',
            r'(?:presentación para|presentation for)[\s:]+(.+)',
        ],
        'industria': [
            r'(?:industria|industry|sector)[\s:]+(.+)',
            r'(?:rubro|giro)[\s:]+(.+)',
        ],
        'problema': [
            r'(?:problema|problemática|issue|problem|desafío|challenge)[\s:]*\n?(.+?)(?=\n\n|\n[A-Z]|objetivo|alcance|$)',
            r'(?:contexto|background|situación actual)[\s:]*\n?(.+?)(?=\n\n|\n[A-Z]|objetivo|$)',
        ],
        'objetivo_general': [
            r'(?:objetivo general|general objective|main objective|objetivo principal)[\s:]*\n?(.+?)(?=\n\n|\n[A-Z]|objetivo secundario|alcance|$)',
            r'(?:objetivo[\s:]+)(.+?)(?=\n\n|\n[A-Z]|$)',
        ],
        'objetivos_secundarios': [
            r'(?:objetivos secundarios|secondary objectives|objetivos específicos|specific objectives)[\s:]*\n?(.+?)(?=\n\n|\n[A-Z]|alcance|$)',
        ],
        'alcance_funcional': [
            r'(?:alcance funcional|functional scope|funcionalidades|features)[\s:]*\n?(.+?)(?=\n\n|\n[A-Z]|alcance técnico|$)',
        ],
        'alcance_tecnico': [
            r'(?:alcance técnico|technical scope|tecnología|technology|stack tecnológico)[\s:]*\n?(.+?)(?=\n\n|\n[A-Z]|alcance geográfico|$)',
        ],
        'alcance_geografico': [
            r'(?:alcance geográfico|geographic scope|cobertura|coverage|ubicación|location)[\s:]*\n?(.+?)(?=\n\n|\n[A-Z]|limitaciones|$)',
        ],
        'limitaciones': [
            r'(?:limitaciones|limitations|exclusiones|exclusions|fuera de alcance|out of scope)[\s:]*\n?(.+?)(?=\n\n|\n[A-Z]|$)',
        ],
    }
    
    # Lista de industrias conocidas para detección
    KNOWN_INDUSTRIES = [
        'banca', 'banking', 'financiero', 'financial', 'seguros', 'insurance',
        'retail', 'comercio', 'minería', 'mining', 'energía', 'energy',
        'telecomunicaciones', 'telecom', 'salud', 'healthcare', 'farmacéutico',
        'pharmaceutical', 'manufactura', 'manufacturing', 'automotriz', 'automotive',
        'tecnología', 'technology', 'gobierno', 'government', 'público', 'public',
        'educación', 'education', 'construcción', 'construction', 'inmobiliario',
        'real estate', 'transporte', 'transportation', 'logística', 'logistics',
        'agricultura', 'agriculture', 'alimentos', 'food', 'bebidas', 'beverages',
        'medios', 'media', 'entretenimiento', 'entertainment', 'turismo', 'tourism',
        'hotelería', 'hospitality', 'consultoria', 'consulting', 'legal',
        'recursos humanos', 'human resources', 'hr', 'utilities', 'servicios públicos'
    ]

    def __init__(self):
        """Inicializa el procesador de documentos"""
        self._check_dependencies()
    
    def _check_dependencies(self) -> Dict[str, bool]:
        """Verifica las dependencias disponibles"""
        return {
            "pdf": PDF_AVAILABLE or PYPDF_AVAILABLE,
            "docx": DOCX_AVAILABLE,
            "pptx": PPTX_AVAILABLE
        }
    
    @staticmethod
    def get_document_type(filename: str) -> DocumentType:
        """Determina el tipo de documento basado en la extensión"""
        ext = Path(filename).suffix.lower()
        type_map = {
            '.pdf': DocumentType.PDF,
            '.docx': DocumentType.WORD,
            '.doc': DocumentType.WORD,
            '.pptx': DocumentType.POWERPOINT,
            '.ppt': DocumentType.POWERPOINT,
        }
        return type_map.get(ext, DocumentType.UNKNOWN)
    
    @staticmethod
    def is_supported_file(filename: str) -> bool:
        """Verifica si el archivo es de un tipo soportado"""
        return DocumentProcessor.get_document_type(filename) != DocumentType.UNKNOWN
    
    def process_document(self, file_path: str) -> ProcessingResult:
        """
        Procesa un documento y extrae información estructurada.
        
        Args:
            file_path: Ruta al archivo a procesar
            
        Returns:
            ProcessingResult con el contenido extraído o error
        """
        filename = os.path.basename(file_path)
        doc_type = self.get_document_type(filename)
        
        if doc_type == DocumentType.UNKNOWN:
            return ProcessingResult(
                filename=filename,
                status=ProcessingStatus.ERROR,
                error_message=f"Tipo de archivo no soportado: {filename}"
            )
        
        try:
            # Extraer texto según el tipo de documento
            if doc_type == DocumentType.PDF:
                raw_text, metadata = self._extract_from_pdf(file_path)
            elif doc_type == DocumentType.WORD:
                raw_text, metadata = self._extract_from_word(file_path)
            elif doc_type == DocumentType.POWERPOINT:
                raw_text, metadata = self._extract_from_pptx(file_path)
            else:
                raise ValueError(f"Tipo de documento no implementado: {doc_type}")
            
            # Extraer información estructurada del texto
            content = self._extract_structured_content(raw_text, metadata)
            
            return ProcessingResult(
                filename=filename,
                status=ProcessingStatus.SUCCESS,
                content=content
            )
            
        except Exception as e:
            return ProcessingResult(
                filename=filename,
                status=ProcessingStatus.ERROR,
                error_message=str(e)
            )
    
    def _extract_from_pdf(self, file_path: str) -> Tuple[str, DocumentMetadata]:
        """Extrae texto y metadatos de un PDF"""
        if not PDF_AVAILABLE and not PYPDF_AVAILABLE:
            raise ImportError("No hay librerías PDF disponibles. Instale pdfplumber o pypdf.")
        
        metadata = DocumentMetadata(
            filename=os.path.basename(file_path),
            file_type="pdf"
        )
        
        text_parts = []
        
        # Usar pdfplumber para extracción de texto (mejor calidad)
        if PDF_AVAILABLE:
            with pdfplumber.open(file_path) as pdf:
                metadata.page_count = len(pdf.pages)
                for page in pdf.pages:
                    page_text = page.extract_text()
                    if page_text:
                        text_parts.append(page_text)
        
        # Obtener metadatos adicionales con pypdf
        if PYPDF_AVAILABLE:
            reader = PdfReader(file_path)
            if reader.metadata:
                metadata.title = reader.metadata.get('/Title', '') or ''
                metadata.author = reader.metadata.get('/Author', '') or ''
                if reader.metadata.get('/CreationDate'):
                    metadata.creation_date = str(reader.metadata.get('/CreationDate', ''))
        
        raw_text = "\n\n".join(text_parts)
        metadata.word_count = len(raw_text.split())
        
        return raw_text, metadata
    
    def _extract_from_word(self, file_path: str) -> Tuple[str, DocumentMetadata]:
        """Extrae texto y metadatos de un documento Word"""
        if not DOCX_AVAILABLE:
            raise ImportError("python-docx no está disponible. Instale con: pip install python-docx")
        
        doc = DocxDocument(file_path)
        
        metadata = DocumentMetadata(
            filename=os.path.basename(file_path),
            file_type="docx"
        )
        
        # Extraer metadatos
        core_props = doc.core_properties
        metadata.title = core_props.title or ""
        metadata.author = core_props.author or ""
        if core_props.created:
            metadata.creation_date = str(core_props.created)
        
        # Extraer texto de párrafos
        text_parts = []
        for para in doc.paragraphs:
            if para.text.strip():
                text_parts.append(para.text)
        
        # Extraer texto de tablas
        for table in doc.tables:
            for row in table.rows:
                row_text = []
                for cell in row.cells:
                    if cell.text.strip():
                        row_text.append(cell.text.strip())
                if row_text:
                    text_parts.append(" | ".join(row_text))
        
        raw_text = "\n\n".join(text_parts)
        metadata.word_count = len(raw_text.split())
        metadata.page_count = max(1, metadata.word_count // 500)  # Estimación
        
        return raw_text, metadata
    
    def _extract_from_pptx(self, file_path: str) -> Tuple[str, DocumentMetadata]:
        """Extrae texto y metadatos de una presentación PowerPoint"""
        if not PPTX_AVAILABLE:
            raise ImportError("python-pptx no está disponible. Instale con: pip install python-pptx")
        
        prs = Presentation(file_path)
        
        metadata = DocumentMetadata(
            filename=os.path.basename(file_path),
            file_type="pptx",
            page_count=len(prs.slides)
        )
        
        # Extraer metadatos
        if prs.core_properties:
            metadata.title = prs.core_properties.title or ""
            metadata.author = prs.core_properties.author or ""
            if prs.core_properties.created:
                metadata.creation_date = str(prs.core_properties.created)
        
        # Extraer texto de cada slide
        text_parts = []
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_text = [f"--- Slide {slide_num} ---"]
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text)
                
                # Extraer texto de tablas
                if shape.has_table:
                    for row in shape.table.rows:
                        row_text = []
                        for cell in row.cells:
                            if cell.text.strip():
                                row_text.append(cell.text.strip())
                        if row_text:
                            slide_text.append(" | ".join(row_text))
            
            if len(slide_text) > 1:  # Más que solo el encabezado del slide
                text_parts.append("\n".join(slide_text))
        
        raw_text = "\n\n".join(text_parts)
        metadata.word_count = len(raw_text.split())
        
        return raw_text, metadata
    
    def _extract_structured_content(self, raw_text: str, metadata: DocumentMetadata) -> ExtractedContent:
        """
        Extrae información estructurada del texto usando patrones y heurísticas.
        """
        content = ExtractedContent(
            metadata=metadata,
            raw_text=raw_text
        )
        
        # Normalizar texto para búsqueda
        text_lower = raw_text.lower()
        
        # Extraer título
        content.titulo_propuesta = self._extract_title(raw_text, metadata)
        
        # Extraer cliente
        content.cliente = self._extract_client(raw_text)
        
        # Extraer industria
        content.industria = self._extract_industry(raw_text)
        
        # Extraer palabras clave
        content.palabras_claves = self._extract_keywords(raw_text)
        
        # Extraer problema/contexto
        content.problema = self._extract_section(raw_text, 'problema')
        
        # Extraer objetivo general
        content.objetivo_general = self._extract_section(raw_text, 'objetivo_general')
        
        # Extraer objetivos secundarios
        objetivos_text = self._extract_section(raw_text, 'objetivos_secundarios')
        content.objetivos_secundarios = self._parse_list_items(objetivos_text)
        
        # Extraer alcances
        content.alcance_funcional = self._extract_section(raw_text, 'alcance_funcional')
        content.alcance_tecnico = self._extract_section(raw_text, 'alcance_tecnico')
        content.alcance_geografico = self._extract_section(raw_text, 'alcance_geografico')
        
        # Extraer limitaciones
        content.limitaciones = self._extract_section(raw_text, 'limitaciones')
        
        # Si no hay limitaciones, inferirlas
        if not content.limitaciones:
            content.limitaciones = self._infer_limitations(content)
        
        return content
    
    def _extract_title(self, text: str, metadata: DocumentMetadata) -> str:
        """Extrae el título de la propuesta"""
        # Primero intentar desde metadatos
        if metadata.title:
            return metadata.title.strip()
        
        # Buscar en el texto con patrones
        for pattern in self.SECTION_PATTERNS['titulo']:
            match = re.search(pattern, text, re.IGNORECASE | re.MULTILINE)
            if match:
                return match.group(1).strip()[:200]
        
        # Usar la primera línea significativa como título
        lines = [l.strip() for l in text.split('\n') if l.strip() and len(l.strip()) > 10]
        if lines:
            # Buscar línea que parezca título (mayúsculas, sin puntuación final)
            for line in lines[:5]:
                if not line.endswith('.') and len(line) < 150:
                    return line
            return lines[0][:200]
        
        return ""
    
    def _extract_client(self, text: str) -> str:
        """Extrae el nombre del cliente"""
        for pattern in self.SECTION_PATTERNS['cliente']:
            match = re.search(pattern, text, re.IGNORECASE)
            if match:
                client = match.group(1).strip()
                # Limpiar y limitar
                client = re.sub(r'\s+', ' ', client)
                return client[:100]
        
        # Buscar patrones comunes de nombres de empresas
        company_patterns = [
            r'([A-Z][A-Za-z]+(?:\s+[A-Z][A-Za-z]+)*\s+(?:S\.?A\.?|SpA|Ltda|LLC|Inc|Corp))',
            r'(?:empresa|company|organización)\s+([A-Z][A-Za-z\s]+)',
        ]
        for pattern in company_patterns:
            match = re.search(pattern, text)
            if match:
                return match.group(1).strip()[:100]
        
        return ""
    
    def _extract_industry(self, text: str) -> str:
        """Extrae la industria del documento"""
        text_lower = text.lower()
        
        # Buscar con patrones específicos
        for pattern in self.SECTION_PATTERNS['industria']:
            match = re.search(pattern, text, re.IGNORECASE)
            if match:
                return match.group(1).strip()[:50]
        
        # Buscar industrias conocidas en el texto
        for industry in self.KNOWN_INDUSTRIES:
            if industry.lower() in text_lower:
                return industry.capitalize()
        
        return ""
    
    def _extract_keywords(self, text: str, max_keywords: int = 15) -> List[str]:
        """Extrae palabras clave del documento"""
        # Tokenizar y limpiar
        text_clean = text.lower()
        text_clean = re.sub(r'[^\w\s]', ' ', text_clean)
        words = text_clean.split()
        
        # Filtrar stopwords y palabras cortas
        filtered_words = [
            w for w in words 
            if w not in self.STOPWORDS_ES 
            and len(w) > 3 
            and not w.isdigit()
        ]
        
        # Contar frecuencias
        word_freq = Counter(filtered_words)
        
        # Obtener las más comunes
        keywords = [word for word, _ in word_freq.most_common(max_keywords)]
        
        return keywords
    
    def _extract_section(self, text: str, section_name: str) -> str:
        """Extrae una sección específica del texto"""
        patterns = self.SECTION_PATTERNS.get(section_name, [])
        
        for pattern in patterns:
            match = re.search(pattern, text, re.IGNORECASE | re.DOTALL)
            if match:
                content = match.group(1).strip()
                # Limpiar y limitar
                content = re.sub(r'\s+', ' ', content)
                return content[:2000]  # Limitar longitud
        
        return ""
    
    def _parse_list_items(self, text: str, max_items: int = 3) -> List[str]:
        """Parsea elementos de una lista desde texto"""
        if not text:
            return []
        
        items = []
        
        # Buscar items con bullets o números
        patterns = [
            r'[-•●○]\s*(.+?)(?=[-•●○]|\n\n|$)',
            r'\d+[.)]\s*(.+?)(?=\d+[.)]|\n\n|$)',
            r'\n\s*(.+?)(?=\n|$)',
        ]
        
        for pattern in patterns:
            matches = re.findall(pattern, text, re.MULTILINE)
            if matches:
                items = [m.strip() for m in matches if m.strip() and len(m.strip()) > 10]
                break
        
        # Si no hay items estructurados, dividir por oraciones
        if not items:
            sentences = re.split(r'[.;]\s+', text)
            items = [s.strip() for s in sentences if s.strip() and len(s.strip()) > 10]
        
        return items[:max_items]
    
    def _infer_limitations(self, content: ExtractedContent) -> str:
        """Infiere limitaciones basado en el contenido extraído"""
        limitations = []
        
        # Basado en el alcance funcional
        if content.alcance_funcional:
            limitations.append(
                "El alcance se limita a las funcionalidades especificadas en el documento."
            )
        
        # Basado en el alcance técnico
        if content.alcance_tecnico:
            limitations.append(
                "Las soluciones técnicas se enmarcan en las tecnologías descritas."
            )
        
        # Basado en el alcance geográfico
        if content.alcance_geografico:
            limitations.append(
                f"La cobertura geográfica está limitada a: {content.alcance_geografico[:100]}."
            )
        
        # Limitaciones genéricas si no hay información específica
        if not limitations:
            limitations = [
                "El proyecto no incluye desarrollo de funcionalidades fuera del alcance definido.",
                "No se consideran integraciones con sistemas no especificados.",
                "El soporte post-implementación está sujeto a acuerdo separado."
            ]
        
        return " ".join(limitations)


class MultiDocumentProcessor:
    """
    Procesador para múltiples documentos con consolidación de información.
    """
    
    def __init__(self):
        self.processor = DocumentProcessor()
        self.results: List[ProcessingResult] = []
    
    def process_documents(self, file_paths: List[str]) -> List[ProcessingResult]:
        """
        Procesa múltiples documentos.
        
        Args:
            file_paths: Lista de rutas a archivos
            
        Returns:
            Lista de resultados de procesamiento
        """
        self.results = []
        
        for file_path in file_paths:
            result = self.processor.process_document(file_path)
            self.results.append(result)
        
        return self.results
    
    def consolidate_content(self) -> ExtractedContent:
        """
        Consolida el contenido de todos los documentos procesados exitosamente,
        eliminando duplicados.
        
        Returns:
            ExtractedContent consolidado
        """
        consolidated = ExtractedContent()
        
        # Recolectar información de todos los documentos exitosos
        all_titles = []
        all_keywords = []
        all_clientes = []
        all_industrias = []
        all_problemas = []
        all_obj_generales = []
        all_obj_secundarios = []
        all_alcance_func = []
        all_alcance_tec = []
        all_alcance_geo = []
        all_limitaciones = []
        all_raw_text = []
        
        for result in self.results:
            if result.status == ProcessingStatus.SUCCESS and result.content:
                content = result.content
                
                if content.titulo_propuesta:
                    all_titles.append(content.titulo_propuesta)
                if content.palabras_claves:
                    all_keywords.extend(content.palabras_claves)
                if content.cliente:
                    all_clientes.append(content.cliente)
                if content.industria:
                    all_industrias.append(content.industria)
                if content.problema:
                    all_problemas.append(content.problema)
                if content.objetivo_general:
                    all_obj_generales.append(content.objetivo_general)
                if content.objetivos_secundarios:
                    all_obj_secundarios.extend(content.objetivos_secundarios)
                if content.alcance_funcional:
                    all_alcance_func.append(content.alcance_funcional)
                if content.alcance_tecnico:
                    all_alcance_tec.append(content.alcance_tecnico)
                if content.alcance_geografico:
                    all_alcance_geo.append(content.alcance_geografico)
                if content.limitaciones:
                    all_limitaciones.append(content.limitaciones)
                if content.raw_text:
                    all_raw_text.append(content.raw_text)
        
        # Consolidar eliminando duplicados
        consolidated.titulo_propuesta = self._get_best_value(all_titles)
        consolidated.palabras_claves = self._deduplicate_list(all_keywords)[:15]
        consolidated.cliente = self._get_best_value(all_clientes)
        consolidated.industria = self._get_best_value(all_industrias)
        consolidated.problema = self._merge_texts(all_problemas)
        consolidated.objetivo_general = self._get_best_value(all_obj_generales)
        consolidated.objetivos_secundarios = self._deduplicate_list(all_obj_secundarios)[:3]
        consolidated.alcance_funcional = self._merge_texts(all_alcance_func)
        consolidated.alcance_tecnico = self._merge_texts(all_alcance_tec)
        consolidated.alcance_geografico = self._get_best_value(all_alcance_geo)
        consolidated.limitaciones = self._merge_texts(all_limitaciones)
        consolidated.raw_text = "\n\n---\n\n".join(all_raw_text)
        
        # Metadata consolidada
        consolidated.metadata = DocumentMetadata(
            filename="consolidated",
            file_type="multiple",
            page_count=sum(
                r.content.metadata.page_count 
                for r in self.results 
                if r.status == ProcessingStatus.SUCCESS and r.content
            ),
            word_count=sum(
                r.content.metadata.word_count 
                for r in self.results 
                if r.status == ProcessingStatus.SUCCESS and r.content
            )
        )
        
        return consolidated
    
    def _get_best_value(self, values: List[str]) -> str:
        """Obtiene el mejor valor de una lista (el más largo no vacío)"""
        valid_values = [v for v in values if v and v.strip()]
        if not valid_values:
            return ""
        # Retornar el más largo
        return max(valid_values, key=len)
    
    def _deduplicate_list(self, items: List[str]) -> List[str]:
        """Elimina duplicados de una lista manteniendo el orden"""
        seen = set()
        result = []
        for item in items:
            item_lower = item.lower().strip()
            if item_lower and item_lower not in seen:
                seen.add(item_lower)
                result.append(item)
        return result
    
    def _merge_texts(self, texts: List[str], max_length: int = 3000) -> str:
        """Combina textos eliminando duplicados y limitando longitud"""
        unique_sentences = set()
        merged_parts = []
        
        for text in texts:
            if not text:
                continue
            # Dividir en oraciones
            sentences = re.split(r'(?<=[.!?])\s+', text)
            for sentence in sentences:
                sentence = sentence.strip()
                sentence_key = sentence.lower()
                if sentence and sentence_key not in unique_sentences:
                    unique_sentences.add(sentence_key)
                    merged_parts.append(sentence)
        
        merged = " ".join(merged_parts)
        if len(merged) > max_length:
            merged = merged[:max_length] + "..."
        
        return merged
    
    def get_processing_summary(self) -> Dict:
        """Obtiene un resumen del procesamiento"""
        success_count = sum(1 for r in self.results if r.status == ProcessingStatus.SUCCESS)
        error_count = sum(1 for r in self.results if r.status == ProcessingStatus.ERROR)
        
        return {
            "total_documents": len(self.results),
            "successful": success_count,
            "failed": error_count,
            "results": [
                {
                    "filename": r.filename,
                    "status": r.status.value,
                    "error": r.error_message if r.status == ProcessingStatus.ERROR else None
                }
                for r in self.results
            ]
        }


# Funciones de utilidad para uso en Streamlit

def save_uploaded_file(uploaded_file, destination_folder: str) -> str:
    """
    Guarda un archivo subido en una carpeta destino.
    
    Args:
        uploaded_file: Archivo de Streamlit (UploadedFile)
        destination_folder: Carpeta donde guardar el archivo
        
    Returns:
        Ruta completa del archivo guardado
    """
    os.makedirs(destination_folder, exist_ok=True)
    file_path = os.path.join(destination_folder, uploaded_file.name)
    
    with open(file_path, "wb") as f:
        f.write(uploaded_file.getbuffer())
    
    return file_path


def cleanup_temp_files(folder: str):
    """Limpia archivos temporales de una carpeta"""
    if os.path.exists(folder):
        for file in os.listdir(folder):
            file_path = os.path.join(folder, file)
            try:
                if os.path.isfile(file_path):
                    os.unlink(file_path)
            except Exception as e:
                print(f"Error al eliminar {file_path}: {e}")


def get_file_icon(filename: str) -> str:
    """Retorna un emoji según el tipo de archivo"""
    doc_type = DocumentProcessor.get_document_type(filename)
    icons = {
        DocumentType.PDF: "📄",
        DocumentType.WORD: "📝",
        DocumentType.POWERPOINT: "📊",
        DocumentType.UNKNOWN: "❓"
    }
    return icons.get(doc_type, "📁")


def format_keywords(keywords: List[str]) -> str:
    """Formatea lista de keywords como string separado por comas"""
    return ", ".join(keywords)


def check_dependencies() -> Dict[str, bool]:
    """
    Verifica qué dependencias están instaladas.
    
    Returns:
        Diccionario con el estado de cada dependencia
    """
    return {
        "pdfplumber": PDF_AVAILABLE,
        "pypdf": PYPDF_AVAILABLE,
        "python-docx": DOCX_AVAILABLE,
        "python-pptx": PPTX_AVAILABLE
    }


def install_dependencies():
    """
    Instala las dependencias necesarias usando pip.
    
    Nota: Esta función está pensada para ser ejecutada manualmente o 
    al inicio de la aplicación si faltan dependencias.
    """
    packages = [
        "pdfplumber",
        "pypdf",
        "python-docx",
        "python-pptx"
    ]
    
    for package in packages:
        try:
            subprocess.check_call([
                "pip", "install", package, "--break-system-packages", "-q"
            ])
        except subprocess.CalledProcessError as e:
            print(f"Error instalando {package}: {e}")


if __name__ == "__main__":
    # Ejemplo de uso
    print("=== Document Processor Module ===")
    print("\nDependencias disponibles:")
    deps = check_dependencies()
    for dep, available in deps.items():
        status = "✅" if available else "❌"
        print(f"  {status} {dep}")
    
    print("\nTipos de archivo soportados:")
    print("  - PDF (.pdf)")
    print("  - Word (.docx)")
    print("  - PowerPoint (.pptx)")